from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 1: Title Slide =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Sprint 2: LLM Benchmark Project"
title_p = title_frame.paragraphs[0]
title_p.font.size = Pt(48)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(52, 73, 94)
title_p.alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Format Comparison Study: CSV vs HTML vs JSON vs XML"
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.font.size = Pt(28)
subtitle_p.font.color.rgb = RGBColor(52, 152, 219)
subtitle_p.alignment = PP_ALIGN.CENTER

# Add name and date
name_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
name_frame = name_box.text_frame
name_frame.text = "Nishit Suthar\nMay 2026"
name_p = name_frame.paragraphs[0]
name_p.font.size = Pt(20)
name_p.font.color.rgb = RGBColor(127, 140, 141)
name_p.alignment = PP_ALIGN.CENTER

# Add key insight box
insight_box = slide1.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(7), Inches(0.8))
insight_frame = insight_box.text_frame
insight_frame.text = "🎯 Major Discovery: All Formats Converge to 80% in Individual Mode"
insight_p = insight_frame.paragraphs[0]
insight_p.font.size = Pt(18)
insight_p.font.bold = True
insight_p.font.color.rgb = RGBColor(231, 76, 60)
insight_p.alignment = PP_ALIGN.CENTER

# ===== SLIDE 2: Methodology & Key Insights =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title2_frame = title2_box.text_frame
title2_frame.text = "Methodology & Key Insights"
title2_p = title2_frame.paragraphs[0]
title2_p.font.size = Pt(36)
title2_p.font.bold = True
title2_p.font.color.rgb = RGBColor(52, 73, 94)
title2_p.alignment = PP_ALIGN.CENTER

# Content box
content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Section 1: Synthetic Data Generation
p1 = content_frame.paragraphs[0]
p1.text = "1. Synthetic Data Generation (Local, NOT LLM-Generated)"
p1.font.size = Pt(22)
p1.font.bold = True
p1.font.color.rgb = RGBColor(41, 128, 185)
p1.space_after = Pt(8)

p2 = content_frame.add_paragraph()
p2.text = "   • Generated 32,550 music industry records using Python Faker library locally"
p2.font.size = Pt(16)
p2.level = 1
p2.space_after = Pt(4)

p3 = content_frame.add_paragraph()
p3.text = "   • 11 relational tables with proper foreign key constraints and referential integrity"
p3.font.size = Pt(16)
p3.level = 1
p3.space_after = Pt(12)

# Section 2: Strategic Data Degradation
p4 = content_frame.add_paragraph()
p4.text = "2. Strategic Data Degradation for Testing"
p4.font.size = Pt(22)
p4.font.bold = True
p4.font.color.rgb = RGBColor(230, 126, 34)
p4.space_after = Pt(8)

p5 = content_frame.add_paragraph()
p5.text = "   • Sampled to 567 records (2.1%) to test LLM behavior with limited data"
p5.font.size = Pt(16)
p5.level = 1
p5.space_after = Pt(4)

p6 = content_frame.add_paragraph()
p6.text = "   • Intentionally kept stale metadata (albums.total_tracks = 12-16 from full dataset)"
p6.font.size = Pt(16)
p6.level = 1
p6.space_after = Pt(4)

p7 = content_frame.add_paragraph()
p7.text = "   • Actual tracks per album after sampling: 0-2 (testing metadata trust)"
p7.font.size = Pt(16)
p7.level = 1
p7.space_after = Pt(12)

# Section 3: Format Conversion
p8 = content_frame.add_paragraph()
p8.text = "3. Multi-Format Conversion & Testing"
p8.font.size = Pt(22)
p8.font.bold = True
p8.font.color.rgb = RGBColor(46, 204, 113)
p8.space_after = Pt(8)

p9 = content_frame.add_paragraph()
p9.text = "   • Converted sampled data to 4 formats: CSV (41KB), HTML (168KB), JSON (173KB), XML (211KB)"
p9.font.size = Pt(16)
p9.level = 1
p9.space_after = Pt(4)

p10 = content_frame.add_paragraph()
p10.text = "   • 20 benchmark questions (6 Medium, 7 Hard, 7 Extremely Hard)"
p10.font.size = Pt(16)
p10.level = 1
p10.space_after = Pt(4)

p11 = content_frame.add_paragraph()
p11.text = "   • Ground truth verified against NeonDB (PostgreSQL)"
p11.font.size = Pt(16)
p11.level = 1

# ===== SLIDE 3: Results & Visualizations =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
title3_frame = title3_box.text_frame
title3_frame.text = "Results: Perfect Format Convergence at 80%"
title3_p = title3_frame.paragraphs[0]
title3_p.font.size = Pt(32)
title3_p.font.bold = True
title3_p.font.color.rgb = RGBColor(52, 73, 94)
title3_p.alignment = PP_ALIGN.CENTER

# Add comprehensive graph
img_path = '/Users/I772947/personal work/LLM Benchmark Team Project/LLM_Benchmark_Team_Project_2026/Sprint 2/comprehensive_benchmark_analysis.png'
slide3.shapes.add_picture(img_path, Inches(0.3), Inches(0.9), width=Inches(9.4))

# Add key insights box at bottom
insights_box = slide3.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.9))
insights_frame = insights_box.text_frame
insights_frame.word_wrap = True

# Key findings
k1 = insights_frame.paragraphs[0]
k1.text = "🔑 Key Findings:"
k1.font.size = Pt(16)
k1.font.bold = True
k1.font.color.rgb = RGBColor(231, 76, 60)

k2 = insights_frame.add_paragraph()
k2.text = "✓ CSV: ONLY format to correctly handle stale metadata (Q007: Jazz = 2.0) | 100% on Hard questions"
k2.font.size = Pt(13)
k2.font.color.rgb = RGBColor(39, 174, 96)
k2.space_after = Pt(2)

k3 = insights_frame.add_paragraph()
k3.text = "✓ XML: Biggest improvement (+15%) from 65% → 80% | HTML/JSON/XML: All used wrong aggregation (Jazz = 0.67)"
k3.font.size = Pt(13)
k3.font.color.rgb = RGBColor(41, 128, 185)
k3.space_after = Pt(2)

k4 = insights_frame.add_paragraph()
k4.text = "✓ Individual mode eliminates format advantages: All converge to 80% ceiling despite different paths"
k4.font.size = Pt(13)
k4.font.color.rgb = RGBColor(142, 68, 173)

# Save presentation
prs.save('/Users/I772947/personal work/LLM Benchmark Team Project/LLM_Benchmark_Team_Project_2026/Sprint 2/Sprint2_Presentation_Nishit_Suthar.pptx')
print("✅ Presentation created: Sprint2_Presentation_Nishit_Suthar.pptx")
